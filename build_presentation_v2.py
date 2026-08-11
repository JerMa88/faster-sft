"""
AT&T CDO Data Science Show
"Bridging the Knowing-Using Gap in Enterprise AI"
RESEARCHER PERSPECTIVE REWRITE

Audience: Product Owners, Investors, Business Personnel
Goal: Educate them on SFT, its technical limitations in data science,
      related works on model architecture analysis, and evolving SFT methods.
      Last slide: Researcher's methodology introduction.

Template: slide_format.pptx
Output:   AT&T_CDO_DataScienceShow_SFT_KUG_v2.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import os

# ──────────────────────────────────────────────────────────────
# Constants
# ──────────────────────────────────────────────────────────────
FOOTER_TEXT = (
    "AT&T CDO Data Science Show / August 2026 / "
    "© 2026 AT&T Intellectual Property – AT&T Proprietary (Internal)"
)
FIG_DIR = "related_works/figures/"

# Image paths
IMG_KUG_CURVE    = FIG_DIR + "crop_kug_learning_curve.png"
IMG_KUG_RIGHT    = FIG_DIR + "crop_kug_right.png"
IMG_KUG_PAGE     = FIG_DIR + "fig_kug_concept_p2.png"
IMG_SELFPATCH_M  = FIG_DIR + "crop_selfpatch_method_comparison.png"
IMG_ALGO         = FIG_DIR + "crop_algorithm1_selfpatch.png"
IMG_HEATMAP_DYN  = FIG_DIR + "crop_fig4_permeation.png"
IMG_CLUSTER      = FIG_DIR + "crop_cluster_heatmaps.png"
IMG_TABLE4       = FIG_DIR + "crop_table4_oracle_results.png"
IMG_OWN_SCATTER  = FIG_DIR + "own_fig1_kug_scatter.png"
IMG_OWN_CURVES   = FIG_DIR + "own_fig2_learning_curves.png"
IMG_OWN_GAINS    = FIG_DIR + "own_fig3_gains.png"
IMG_OWN_ORACLE   = FIG_DIR + "own_fig5_oracle_headroom.png"

# ──────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────

def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout not found: {name!r}")


def add_slide(prs, layout_name):
    return prs.slides.add_slide(get_layout(prs, layout_name))


def set_ph(slide, idx, text, bold=False, font_size=None, italic=False, color=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            if bold:      run.font.bold = True
            if italic:    run.font.italic = True
            if font_size: run.font.size = Pt(font_size)
            if color:     run.font.color.rgb = color
            return ph
    return None


def set_ph_paras(slide, idx, paras, default_size=None):
    """
    paras: list of dicts with keys: text, bold, italic, size, color, bullet
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, para in enumerate(paras):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = para.get("text", "")
                if para.get("bold"):    run.font.bold = True
                if para.get("italic"):  run.font.italic = True
                sz = para.get("size", default_size)
                if sz:                  run.font.size = Pt(sz)
                col = para.get("color")
                if col:                 run.font.color.rgb = col
            return ph
    return None


def set_footer(slide):
    set_ph(slide, 12, FOOTER_TEXT)


def add_image(slide, img_path, left, top, width=None, height=None):
    """Add image to slide at given position (in Inches)."""
    if not os.path.exists(img_path):
        print(f"  WARNING: image not found: {img_path}")
        return None
    kw = {}
    if width:  kw["width"]  = Inches(width)
    if height: kw["height"] = Inches(height)
    return slide.shapes.add_picture(img_path, Inches(left), Inches(top), **kw)


def add_textbox(slide, left, top, width, height, text,
                bold=False, italic=False, font_size=12,
                color=None, bg_color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(font_size)
    if color:
        run.font.color.rgb = color
    return txBox


def add_label(slide, left, top, width, height, text, font_size=9,
              bold=False, italic=False, color=None):
    """Small caption/label text box."""
    return add_textbox(slide, left, top, width, height, text,
                       font_size=font_size, bold=bold, italic=italic,
                       color=color or RGBColor(0x44, 0x44, 0x44),
                       align=PP_ALIGN.CENTER, wrap=True)


# ──────────────────────────────────────────────────────────────
# Load template & clear slides
# ──────────────────────────────────────────────────────────────
prs = Presentation("slide_format.pptx")

while len(prs.slides) > 0:
    slide_elem = prs.slides._sldIdLst[0]
    rId = slide_elem.get(qn("r:id"))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
    del prs.slides._sldIdLst[0]

print(f"Template loaded. Layouts: {len(prs.slide_layouts)}")

# Slide dimensions (13.33" x 7.50")
SW = 13.333
SH = 7.5

ATT_BLUE = RGBColor(0x00, 0x9F, 0xDB)
DARK     = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE   = RGBColor(0xE8, 0x6A, 0x00)
GREEN    = RGBColor(0x2E, 0x86, 0x48)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Cover: Title slide
# ══════════════════════════════════════════════════════════════
s1 = add_slide(prs, "Cover w/ Subtitle (White)")
set_ph(s1, 0,
    "Bridging the Knowing-Using Gap in Enterprise AI",
    bold=True)
set_ph(s1, 13,
    "Understanding SFT, Its Technical Limitations, and the Path Forward")
set_ph(s1, 10, "Zerui (Jerry) Ma  ·  AT&T Chief Data Office  ·  AI Research")
set_ph(s1, 11, "CDO Data Science Show  |  August 2026")
print("✓ Slide 1: Cover")


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Real-World Example: AT&T OTel 2.0
# Layout: Title + Subtitle (1/2 Image on Right)
# Inserted as intro hook: why SFT matters in production TODAY
# ══════════════════════════════════════════════════════════════
IMG_OTEL = FIG_DIR + "otel_abstract.png"

s_otel = add_slide(prs, "Title + Subtitle (1/2 Image on Right)")
set_ph(s_otel, 0, "Real Impact: AT&T OTel 2.0 — #1 on Open Telco AI Leaderboard")
set_ph(s_otel, 13,
    "Friday, July 24 2026 · GSMA Open-Telco.ai — the largest and best-performing open-source model built for telecoms")
set_footer(s_otel)

set_ph_paras(s_otel, 18, [
    {"text": "What was built:", "size": 13, "bold": True},
    {"text": "A post-trained version of Gemma 4 31B-IT, trained on 400 billion telecom-specific tokens selected from 1+ trillion processed tokens. Collaboration: AT&T, GSMA, Red Hat, Dell, Microsoft Azure, AMD.", "size": 12},
    {"text": ""},
    {"text": "Why domain-adapted models win:", "size": 13, "bold": True},
    {"text": "The top 3 performers on Open Telco AI benchmarks are ALL domain-adapted models — NOT general-purpose frontier models (GPT-4, Claude, Gemini). OTel 2.0 sits at #1.", "size": 12},
    {"text": ""},
    {"text": "The core limitation that motivated this:", "size": 13, "bold": True},
    {"text": "General-purpose models weren't built with telecoms in mind. Ask one to interpret a 3GPP standard or troubleshoot a live network issue, and the cracks show — not because they lack capability, but because the training data barely touches this domain.", "size": 12, "italic": True},
    {"text": ""},
    {"text": "This is fine-tuning (SFT) delivering measurable competitive advantage.", "size": 13, "bold": True},
    {"text": "The remaining slides explain HOW this works technically, and WHAT breaks when it isn't done correctly.", "size": 12},
])

# Add the OTel abstract image on the right
add_image(s_otel, IMG_OTEL, left=7.15, top=1.3, width=5.95)
add_label(s_otel, 7.15, 6.15, 5.95, 0.4,
          "Source: GSMA / open-telco.ai · July 24, 2026",
          font_size=8, italic=True)
print("✓ Slide 2: OTel 2.0 real-world example")

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — What is an LLM and How Does Fine-Tuning Work?
# Layout: Title + Subtitle (2 column)
# ══════════════════════════════════════════════════════════════

s2 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s2, 0, "What is Supervised Fine-Tuning (SFT)?")
set_ph(s2, 13,
    "The technique that adapts a general-purpose LLM to a specific knowledge domain")
set_footer(s2)

set_ph(s2, 14, "Step 1 — Pre-Training (Once, by vendors)")
set_ph_paras(s2, 18, [
    {"text": "A Large Language Model (LLM) like LLaMA, Qwen, or Gemma is trained on trillions of tokens of internet text. It learns language, grammar, general world knowledge, and reasoning patterns.", "size": 13},
    {"text": ""},
    {"text": "This costs millions of dollars and takes months. AT&T does not do this.", "size": 12, "italic": True},
    {"text": ""},
    {"text": "The result: a model that can write, reason, and answer questions — but knows nothing specific about AT&T.", "size": 13},
])

set_ph(s2, 15, "Step 2 — Fine-Tuning (SFT, done by AT&T)")
set_ph_paras(s2, 19, [
    {"text": "We take the pre-trained model and train it further on AT&T-specific question-answer pairs:", "size": 13},
    {"text": ""},
    {"text": "  Q: Which protein is encoded by gene BRCA1?\n  A: BRCA1 protein (tumor suppressor)", "size": 12, "italic": True},
    {"text": ""},
    {"text": "Training signal: Next-token cross-entropy loss — the model learns to predict the answer token-by-token.", "size": 12},
    {"text": ""},
    {"text": "Using LoRA (Low-Rank Adaptation): instead of updating all ~1-7 billion parameters, we add thin trainable adapter matrices. Fast, cheap, memory-efficient.", "size": 13, "bold": True},
])
print("✓ Slide 2: What is SFT?")

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — The Taxonomy of Knowledge Tasks in SFT
# Layout: Title + Subtitle (2 column)
# ══════════════════════════════════════════════════════════════
s3 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s3, 0, "Two Fundamentally Different Types of Knowledge Tasks")
set_ph(s3, 13,
    "SFT is evaluated on both, but only trained on one — this asymmetry is the root of the problem")
set_footer(s3)

set_ph(s3, 14, "Memorization Task (P_mem) — What SFT Trains On")
set_ph_paras(s3, 18, [
    {"text": "Direct fact recall: given a prompt, produce a stored answer.", "size": 13},
    {"text": ""},
    {"text": "Example (STaRK-Prime biomedical KG):", "size": 12},
    {"text": "  Q: Which drug is transported by gene ABCC11?\n  A: Probenecid", "size": 12, "italic": True},
    {"text": ""},
    {"text": "This is a single-hop lookup. The answer is directly written in the training data. The model memorizes it.", "size": 12},
    {"text": ""},
    {"text": "Accuracy metric: A_mem — does the model produce the correct answer when directly asked?", "size": 12, "bold": True},
])

set_ph(s3, 15, "Generalization Task (P_gen) — What SFT Is Tested On")
set_ph_paras(s3, 19, [
    {"text": "Multi-hop reasoning: two-step inference using injected facts.", "size": 13},
    {"text": ""},
    {"text": "Example (2-hop chaining):", "size": 12},
    {"text": "  Q: What is a side effect of the drug transported by ABCC11?\n  A: Vertigo", "size": 12, "italic": True},
    {"text": ""},
    {"text": "Step 1: ABCC11 → Probenecid  (injected via SFT)\nStep 2: Probenecid → Vertigo  (pre-existing knowledge)", "size": 12},
    {"text": ""},
    {"text": "Accuracy metric: A_gen — can the model chain two facts together to answer a question it wasn't explicitly trained on?", "size": 12, "bold": True},
])
print("✓ Slide 3: Task taxonomy")

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Section Divider: The Knowing-Using Gap
# ══════════════════════════════════════════════════════════════
s4 = add_slide(prs, "Divider 01")
set_ph(s4, 0, "The Knowing-Using Gap (KUG)\n\nWhy SFT Fails at Multi-Hop Reasoning")
set_footer(s4)
print("✓ Slide 4: Divider — KUG")

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — The KUG: Empirical Observation (Figure 1 from paper)
# Layout: Title + Subtitle (1/2 Image on Right) → use blank + manual placement
# ══════════════════════════════════════════════════════════════
s5 = add_slide(prs, "Title + Subtitle (1/2 Image on Right)")
set_ph(s5, 0, "The Knowing-Using Gap: Empirically Observed in SFT")
set_ph(s5, 13, 'Dai et al., 2025 — "Towards Mechanistically Understanding Why Memorized Knowledge Fails to Generalize in LLM Finetuning" (arXiv:2607.08393)')
set_footer(s5)

set_ph_paras(s5, 18, [
    {"text": "After SFT on 1,000 facts across Qwen/LLaMA families:", "size": 13, "bold": True},
    {"text": ""},
    {"text": "• Memorization (A_mem) saturates near 1.0 by epoch ~10.", "size": 12},
    {"text": "• Chaining accuracy (A_gen) rises to only 0.078–0.20 at epoch 30.", "size": 12},
    {"text": "• Gradient vanishes at memorization saturation — no further signal to drive generalization.", "size": 12},
    {"text": ""},
    {"text": "Two symptoms define the KUG:", "size": 13, "bold": True},
    {"text": "  (1) Accuracy Gap: A_mem >> A_gen at convergence.", "size": 12},
    {"text": "  (2) Time Lag: generalization saturates 4–5 epochs later than memorization.", "size": 12},
    {"text": ""},
    {"text": "Key finding: the bottleneck is NOT data scarcity, model size, or training duration.", "size": 12, "bold": True, "italic": True},
])

# Add the figure
add_image(s5, IMG_KUG_CURVE, left=7.4, top=1.5, width=5.7)
add_label(s5, 7.4, 6.2, 5.7, 0.5,
          "Fig. 1(a) from Dai et al. (2025, arXiv:2607.08393)", font_size=8, italic=True)
print("✓ Slide 5: KUG empirical observation")

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — KUG Scale: Confirmed across Models & Domains
# Layout: Title + Subtitle (2 column)
# ══════════════════════════════════════════════════════════════
s6 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s6, 0, "KUG Is Universal: Confirmed Across Model Families & Scales")
set_ph(s6, 13,
    "Dai et al. (2025) tested 6 models (1.5B–8B) on STaRK-Prime (biomedical) and STaRK-MAG (academic). The gap appears in every case.")
set_footer(s6)

set_ph(s6, 14, "STaRK-Prime & STaRK-MAG Benchmark Results (Dai et al., 2025)")
set_ph_paras(s6, 18, [
    {"text": "From Table 3 (arXiv:2607.08393):", "size": 12, "bold": True},
    {"text": ""},
    {"text": "Model          A_mem    A_gen (Chain)    Gap", "size": 11, "italic": True},
    {"text": "Qwen2.5-1.5B   0.998       0.078         0.920", "size": 11},
    {"text": "Qwen2.5-3B     0.997       0.114         0.883", "size": 11},
    {"text": "Qwen2.5-7B     0.996       0.124         0.872", "size": 11},
    {"text": "LLaMA-3.2-1B   0.994       0.102         0.892", "size": 11},
    {"text": "LLaMA-3.2-3B   0.993       0.126         0.867", "size": 11},
    {"text": "LLaMA-3.1-8B   0.986       0.182         0.804", "size": 11},
    {"text": ""},
    {"text": "Interpretation: models memorize near-perfectly (A_mem > 0.98) but generalize poorly (A_gen < 0.20) across all scales.", "size": 12, "bold": True},
])

# Right: our replication data
set_ph(s6, 15, "AT&T CDO Replication (Jerry Ma et al., 2026 — This Work)")
add_image(s6, IMG_OWN_SCATTER, left=7.0, top=1.6, width=6.1)
add_label(s6, 7.0, 6.0, 6.1, 0.5,
          "KUG scatter across 4 models × 2 datasets. All points fall below Amem=Agen diagonal.",
          font_size=8, italic=True)
print("✓ Slide 6: KUG universality")

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Section Divider: Mechanistic Root Cause
# ══════════════════════════════════════════════════════════════
s7 = add_slide(prs, "Divider 01")
set_ph(s7, 0, "Mechanistic Root Cause\n\nWhy Does the KUG Exist?\nThe Knowledge-Circuit Misalignment Hypothesis")
set_footer(s7)
print("✓ Slide 7: Divider — Mechanistic")

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Layer Specialization in Transformers
# Layout: Title + Subtitle (2 column)
# ══════════════════════════════════════════════════════════════
s8 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s8, 0, "Transformer Layer Specialization: Storage vs. Reasoning")
set_ph(s8, 13,
    "Mechanistic interpretability research (Meng et al. 2022 [ROME]; Allen-Zhu & Li 2023 [Physics of LM]) reveals distinct functional zones in LLMs")
set_footer(s8)

set_ph(s8, 14, "Early & Late Layers — Knowledge Storage (MLP neurons)")
set_ph_paras(s8, 18, [
    {"text": "ROME (Meng et al., 2022, NeurIPS): causal tracing identifies that new factual associations are concentrated in MLP layers at specific depths (e.g., layers 13–17 in GPT-2-XL).", "size": 12},
    {"text": ""},
    {"text": "Physics of LM Part 3.1 (Allen-Zhu & Li, 2023): knowledge is only extractable if stored LINEARLY on entity-name token positions. Without augmentation, facts disperse across positions and become inaccessible.", "size": 12},
    {"text": ""},
    {"text": "SFT gradient descent takes the path of least resistance: it satisfies the next-token prediction loss by writing facts into storage layers and STOPS.", "size": 13, "bold": True},
])

set_ph(s8, 15, "Middle Layers — Reasoning Circuits (Attention heads)")
set_ph_paras(s8, 19, [
    {"text": "Multi-hop reasoning (chaining, intersection) executes via attention circuits in the middle ~30–70% of layer depth.", "size": 12},
    {"text": ""},
    {"text": "Physics of LM Part 3.2 (Allen-Zhu & Li, 2023): models excel at retrieval but fail at classification, comparison, and multi-hop tasks without Chain-of-Thought — because CoT explicitly routes facts into the reasoning window.", "size": 12},
    {"text": ""},
    {"text": "The key insight: SFT writes facts into early/late storage, but the gradient VANISHES before it ever nudges the middle-layer reasoning circuits to incorporate the new knowledge.", "size": 13, "bold": True, "italic": True},
    {"text": ""},
    {"text": "The fact is stored — but in the wrong place for the computation that needs it.", "size": 12, "bold": True},
])
print("✓ Slide 8: Layer specialization")

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Self-Patching: The Diagnostic Proof
# Layout: image-heavy, use 1/2 image on right
# ══════════════════════════════════════════════════════════════
s9 = add_slide(prs, "Title + Subtitle (1/2 Image on Right)")
set_ph(s9, 0, "Self-Patching: Causal Proof of the Routing Failure")
set_ph(s9, 13,
    "Dai et al. (2025): a causal intervention that manually routes hidden states from storage layers into reasoning layers")
set_footer(s9)

set_ph_paras(s9, 18, [
    {"text": "Algorithm: Self-Patching Scan (Fig. 3c from Dai et al.):", "size": 13, "bold": True},
    {"text": ""},
    {"text": "For each layer pair (l_src, l_tgt):", "size": 12},
    {"text": "  1. Run memorization prompt P_mem through model", "size": 12},
    {"text": "  2. Cache hidden state at entity position from layer l_src", "size": 12},
    {"text": "  3. Run generalization prompt P_gen through model", "size": 12},
    {"text": "  4. At layer l_tgt, replace entity-position hidden state with cached z", "size": 12},
    {"text": "  5. Measure ΔAcc = Acc(patched) - Acc(unpatched)", "size": 12},
    {"text": ""},
    {"text": "Result (Table 4):", "size": 13, "bold": True},
    {"text": "  w/o patching:  A_gen = 0.078 (chaining, Qwen2.5-1.5B)", "size": 12},
    {"text": "  w/ self-patch:  A_gen = 0.440  →  5.6× improvement!", "size": 12, "bold": True},
    {"text": ""},
    {"text": "Crucial: self-patching introduces no new information. It only RELOCATES existing representations. This proves the knowledge was there all along — just in the wrong layer.", "size": 12, "italic": True, "bold": True},
])

add_image(s9, IMG_SELFPATCH_M, left=7.2, top=1.4, width=5.9)
add_label(s9, 7.2, 5.95, 5.9, 0.45,
          "Fig. 3 (Dai et al., 2025): method comparison — causal tracing, PatchScope, self-patching",
          font_size=8, italic=True)
print("✓ Slide 9: Self-patching diagnostic")

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — The Permeation Heatmap: Knowledge Dynamics During Training
# ══════════════════════════════════════════════════════════════
s10 = add_slide(prs, "Title + Subtitle (1/2 Image on Right)")
set_ph(s10, 0, "Permeation Dynamics: How Knowledge Moves Through Layers During Fine-Tuning")
set_ph(s10, 13,
    "Fig. 4 (Dai et al., 2025) — Each cell (l_src, l_tgt) shows ΔAcc if representation is relocated from source to target layer at that training epoch")
set_footer(s10)

set_ph_paras(s10, 18, [
    {"text": "Reading the heatmap (Fig. 4):", "size": 13, "bold": True},
    {"text": ""},
    {"text": "RED = patching helps (knowledge accessible at l_src, needed at l_tgt)", "size": 12, "color": RGBColor(0xCC, 0x00, 0x00)},
    {"text": "BLUE = patching does not help", "size": 12, "color": RGBColor(0x00, 0x5F, 0xAA)},
    {"text": ""},
    {"text": "Failure case (top row):", "size": 12, "bold": True},
    {"text": "• Off-diagonal red region appears once the fact is memorized, but halts before the diagonal — the middle-layer reasoning circuit never naturally receives the representation.", "size": 12},
    {"text": ""},
    {"text": "Success case (bottom row):", "size": 12, "bold": True},
    {"text": "• Red region gradually expands to cover the diagonal — the fact slowly diffuses into reasoning layers over many epochs.", "size": 12},
    {"text": ""},
    {"text": "This is the visualization of the Knowing-Using Gap at the individual fact level.", "size": 12, "bold": True, "italic": True},
])

add_image(s10, IMG_HEATMAP_DYN, left=7.0, top=1.2, width=6.1)
add_label(s10, 7.0, 5.9, 6.1, 0.5,
          "Fig. 4 (Dai et al., 2025): permeation heatmap of knowledge routing across training epochs",
          font_size=8, italic=True)
print("✓ Slide 10: Permeation heatmap")

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Two-Cluster Structure: The Actionable Finding
# ══════════════════════════════════════════════════════════════
s11 = add_slide(prs, "Title + Subtitle (1/2 Image on Right)")
set_ph(s11, 0, "Effective Patch Locations Concentrate Into Two Clusters")
set_ph(s11, 13,
    "Fig. 5 (Dai et al., 2025) — Consistent pattern across all 6 model families and both knowledge domains")
set_footer(s11)

set_ph_paras(s11, 18, [
    {"text": "Figure 5 reveals a clear structural pattern:", "size": 13, "bold": True},
    {"text": ""},
    {"text": "Cluster 1 (orange): Source in late layers (~0.8L) → target in middle layers (~0.5L)", "size": 12},
    {"text": "Cluster 2 (orange): Source in early layers (~0.1L) → target in middle layers (~0.5L)", "size": 12},
    {"text": ""},
    {"text": "This tells us:", "size": 13, "bold": True},
    {"text": "• Facts are stored in BOTH early AND late layers after SFT.", "size": 12},
    {"text": "• The reasoning bottleneck is consistently at the MIDDLE layers (~L/3 to 2L/3).", "size": 12},
    {"text": "• This two-cluster structure is reproducible across Qwen (1.5B–7B) and LLaMA (1B–8B).", "size": 12},
    {"text": ""},
    {"text": "Fixed heuristic result: using just these two predetermined layer pairs (no per-instance search), the authors recover 58–75% of oracle headroom.", "size": 12, "bold": True},
    {"text": ""},
    {"text": "This is NOT a model artifact — it is a structural property of Transformer fine-tuning.", "size": 12, "italic": True},
])

add_image(s11, IMG_CLUSTER, left=7.0, top=1.8, width=6.1)
add_label(s11, 7.0, 5.9, 6.1, 0.5,
          "Fig. 5 (Dai et al., 2025): effective patch location clusters across 6 models",
          font_size=8, italic=True)
print("✓ Slide 11: Two-cluster structure")

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Section Divider: Related Works
# ══════════════════════════════════════════════════════════════
s12 = add_slide(prs, "Divider 01")
set_ph(s12, 0, "Related Works\n\nWhat the Research Community Has Tried\nand Why Each Approach Falls Short")
set_footer(s12)
print("✓ Slide 12: Divider — Related Works")

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Related Work: Knowledge Editing (ROME, MEMIT, ACE)
# Layout: 2-column
# ══════════════════════════════════════════════════════════════
s13 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s13, 0, "Related Work I: Knowledge Editing (ROME, MEMIT, ACE)")
set_ph(s13, 13,
    "Post-hoc surgical weight edits to inject facts — first to identify the layer-routing problem, but limited to single hops")
set_footer(s13)

set_ph(s13, 14, "ROME & MEMIT (Meng et al., 2022–2023)")
set_ph_paras(s13, 18, [
    {"text": "Causal tracing identifies which MLP layers store a given fact. A rank-one weight update surgically inserts a new fact at that layer.", "size": 12},
    {"text": ""},
    {"text": "✓ Works perfectly for single-hop recall: 'Who is the CEO of OpenAI?' → Sam Altman", "size": 12},
    {"text": ""},
    {"text": "✗ FAILS on 2-hop queries:", "size": 12, "bold": True},
    {"text": "  Cohen et al. (2024, arXiv:2601.04600) — 'On the Limitations of Rank-One Model Editing in Answering Multi-hop Questions': the 'hopping-too-late' problem — facts edited into late layers arrive after the middle-layer composition window has passed.", "size": 11, "italic": True},
    {"text": ""},
    {"text": "ACE (Tang et al., 2024, arXiv:2510.07896):", "size": 12, "bold": True},
    {"text": "Discovers that 'implicit subjects sequentially activate value neurons across transformer layers' — multi-hop reasoning is a layer-sequential chain. Proposes editing at multiple layers, not just one. Still post-hoc, no training signal.", "size": 12},
])

set_ph(s13, 15, "Key Limitation: Post-Hoc, Not Training-Time")
set_ph_paras(s13, 19, [
    {"text": "All knowledge editing methods:", "size": 12, "bold": True},
    {"text": "• Modify weights after training (no gradient flow)", "size": 12},
    {"text": "• Cannot be applied to the LoRA adapters we care about", "size": 12},
    {"text": "• Do not generalize — each fact requires a separate edit operation", "size": 12},
    {"text": "• Fail for batch injection of thousands of enterprise facts", "size": 12},
    {"text": ""},
    {"text": "They independently confirm: layer routing is the problem.", "size": 13, "bold": True},
    {"text": "But none of them propose a training-time loss as the solution.", "size": 13, "italic": True, "bold": True},
    {"text": ""},
    {"text": "These papers serve as independent mechanistic confirmation of our hypothesis from a completely different research direction.", "size": 12},
])
print("✓ Slide 13: Knowledge editing related work")

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — Related Work: Layer-Wise Knowledge Distillation (TinyBERT, OPRD)
# ══════════════════════════════════════════════════════════════
s14 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s14, 0, "Related Work II: Layer-Wise Knowledge Distillation")
set_ph(s14, 13,
    "Representation matching between teacher and student models — technically closest family, but with a fundamental setup difference")
set_footer(s14)

set_ph(s14, 14, "TinyBERT / PKD / OPRD")
set_ph_paras(s14, 18, [
    {"text": "TinyBERT (Jiao et al., 2019, arXiv:1909.10351):", "size": 12, "bold": True},
    {"text": "Distills large BERT → small BERT by MSE-matching hidden states layer-by-layer. Objective: compression, not routing.", "size": 12},
    {"text": ""},
    {"text": "PKD (Sun et al., 2019):", "size": 12, "bold": True},
    {"text": "Patient Knowledge Distillation — match last k intermediate layers of teacher to student. The 'patience' concept is philosophically analogous to our warmup period.", "size": 12},
    {"text": ""},
    {"text": "OPRD (Yang et al., 2026, arXiv:2606.06021):", "size": 12, "bold": True},
    {"text": "On-Policy Representation Distillation — aligns hidden states (not logits) during on-policy rollouts. Proves cosine distance alignment provides lower-variance gradients than KL divergence on output distributions.", "size": 12},
    {"text": ""},
    {"text": "This validates our choice of cosine similarity as the loss metric.", "size": 12, "italic": True, "bold": True},
])

set_ph(s14, 15, "How We Differ: Single-Model, Cross-Prompt")
set_ph_paras(s14, 19, [
    {"text": "Feature Comparison:", "size": 13, "bold": True},
    {"text": ""},
    {"text": "                    Distillation     Ours (AA-SFT)", "size": 11, "italic": True},
    {"text": "Models involved:    Two (T→S)        One (self)", "size": 11},
    {"text": "Prompt used:        Same prompt      P_mem ≠ P_gen", "size": 11},
    {"text": "Alignment target:   Compression      Routing repair", "size": 11},
    {"text": "Teacher in memory:  Yes (large LLM)  No (cached vector)", "size": 11},
    {"text": "Training overhead:  2× model memory  ~1.1× memory", "size": 11},
    {"text": ""},
    {"text": "Our setup is intra-model, cross-prompt: the 'teacher' is the same model reading a different prompt (P_mem), and the 'student' is the same model reading the reasoning prompt (P_gen).", "size": 12, "bold": True},
    {"text": ""},
    {"text": "No prior work has done this combination: (One-model ✓) + (Hidden state loss ✓) + (Cross-prompt ✓) + (Training-time ✓)", "size": 12, "italic": True},
])
print("✓ Slide 14: Distillation related work")

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — Related Work: Activation Steering & MechELK
# ══════════════════════════════════════════════════════════════
s15 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s15, 0, "Related Work III: Activation Steering & Latent Knowledge Elicitation")
set_ph(s15, 13,
    "Inference-time interventions and pre-existing knowledge amplification — different problem, overlapping tools")
set_footer(s15)

set_ph(s15, 14, "Activation Steering (2023–2026)")
set_ph_paras(s15, 18, [
    {"text": "A family of methods that identify directions in the residual stream corresponding to specific behaviors (e.g., 'hallucination', 'tool use'), then steer activations at inference time.", "size": 12},
    {"text": ""},
    {"text": "Representative works: Zou et al. (2023, arXiv: Representation Engineering); ASA (arXiv:2602.02935); AAC (arXiv:2603.10195); FairSteer (arXiv:2504.14492)", "size": 11, "italic": True},
    {"text": ""},
    {"text": "Key limitation: ALL are INFERENCE-TIME only. They patch the model during generation but do NOT update weights. The model does not learn — it must be patched every time.", "size": 12, "bold": True},
    {"text": ""},
    {"text": "What we borrow: the linear probe methodology these papers use to identify 'knowledge-accessible' layers is exactly our Metric 1 (Layer Profiling) and Variant 2 (Probe Loss).", "size": 12, "italic": True},
])

set_ph(s15, 15, "MechELK (arXiv:2605.28825, 2025)")
set_ph_paras(s15, 19, [
    {"text": "MechELK (Mechanistic Elicitation of Latent Knowledge): a 3-stage framework to elicit knowledge that is already encoded in the model's representations but suppressed in outputs.", "size": 12},
    {"text": ""},
    {"text": "Stage 2 is closest to us: 'Amplify via targeted fine-tuning' — they fine-tune targeting specific layers to make EXISTING knowledge more accessible.", "size": 12, "bold": True},
    {"text": ""},
    {"text": "Critical distinction: MechELK targets PRE-EXISTING latent knowledge that the model already has from pre-training. We target NEWLY INJECTED SFT facts that the model did not have before fine-tuning.", "size": 13, "bold": True, "italic": True},
    {"text": ""},
    {"text": "This is a conceptually important difference: eliciting suppressed knowledge vs. routing newly injected knowledge.", "size": 12},
])
print("✓ Slide 15: Activation steering / MechELK")

# ══════════════════════════════════════════════════════════════
# SLIDE 16 — Related Work: Comparison Table (from paper)
# ══════════════════════════════════════════════════════════════
s16 = add_slide(prs, "Title + Subtitle (1/2 Image on Right)")
set_ph(s16, 0, "Landscape Summary: Where AA-SFT Sits Among Prior Work")
set_ph(s16, 13,
    "The combination of five properties is unique to our method across all surveyed literature")
set_footer(s16)

set_ph_paras(s16, 18, [
    {"text": "Five key dimensions:", "size": 13, "bold": True},
    {"text": ""},
    {"text": "(1) Single-model (no teacher needed)", "size": 12},
    {"text": "(2) Hidden-state loss (not output/logit space)", "size": 12},
    {"text": "(3) Cross-prompt (P_mem ≠ P_gen)", "size": 12},
    {"text": "(4) Layer-targeted (profiled, not arbitrary)", "size": 12},
    {"text": "(5) Training-time (weights updated, not inference patched)", "size": 12},
    {"text": ""},
    {"text": "Prior works hit at most 3 of these 5 simultaneously.", "size": 13, "bold": True},
    {"text": ""},
    {"text": "Our novelty statement (from paper §8):", "size": 12},
    {"text": '"To our knowledge, we are the first to propose using intra-model, cross-prompt representation alignment as a training-time auxiliary loss during SFT, motivated by the mechanistically-identified knowledge-circuit misalignment."', "size": 11, "italic": True},
    {"text": ""},
    {"text": "This is a novel application + combination — not a claim of entirely new primitives.", "size": 12, "italic": True},
])

# Add table image (we'll use the results table page from the paper which has our comparison table)
add_image(s16, IMG_OWN_GAINS, left=7.0, top=1.8, width=6.1)
add_label(s16, 7.0, 5.85, 6.1, 0.5,
          "AA-SFT vs. SFT baseline across 4 models × 2 datasets (Ma et al., 2026)",
          font_size=8, italic=True)
print("✓ Slide 16: Related work comparison table")

# ══════════════════════════════════════════════════════════════
# SLIDE 17 — Section Divider: Our Methodology
# ══════════════════════════════════════════════════════════════
s17 = add_slide(prs, "Divider 01")
set_ph(s17, 0, "Our Research Initiative\n\nAlignment-Aware SFT (AA-SFT)\nMethodology")
set_footer(s17)
print("✓ Slide 17: Divider — Our Work")

# ══════════════════════════════════════════════════════════════
# SLIDE 18 [LAST] — AA-SFT Methodology (no results)
# Layout: 2-column
# ══════════════════════════════════════════════════════════════
s18 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(s18, 0, "Alignment-Aware SFT (AA-SFT): Turning the Diagnostic Into a Training Objective")
set_ph(s18, 13,
    "Ma et al. (2026, AT&T CDO) — Methodology only; see paper/main.tex for full results")
set_footer(s18)

set_ph(s18, 14, "The Dual-Forward Pass Training Framework")
set_ph_paras(s18, 18, [
    {"text": "Training objective:", "size": 13, "bold": True},
    {"text": "  L_total = L_SFT(P_mem) + λ · L_align      (λ = 0.1)", "size": 12, "italic": True},
    {"text": ""},
    {"text": "For each fact in the batch:", "size": 12, "bold": True},
    {"text": "  Pass 1 (no_grad): forward P_mem → cache h_E^{l_s} at storage layers l_s^early and l_s^late", "size": 12},
    {"text": "  Pass 2 (grad):    forward P_gen → compute L_SFT + L_align using h_E^{l_t} at reasoning layer l_t", "size": 12},
    {"text": ""},
    {"text": "Layer selection (empirical profiling, not heuristic):", "size": 12, "bold": True},
    {"text": "  l_s^early: first layer where linear probe accuracy > 0.6 (PoLM 3.1 metric)", "size": 12},
    {"text": "  l_s^late:  last layer where probe accuracy > 0.85", "size": 12},
    {"text": "  l_t:       argmax of self-patching gain map A[l_src, l_tgt]", "size": 12},
    {"text": ""},
    {"text": "Warmup: L_align activates only after K=3 epochs (memorization must saturate first)", "size": 12, "italic": True},
    {"text": "At deployment: standard single forward pass — zero inference overhead.", "size": 12, "bold": True},
])

set_ph(s18, 15, "Four Alignment Loss Variants")
set_ph_paras(s18, 19, [
    {"text": "RepDist (Representation Distillation):", "size": 12, "bold": True},
    {"text": "  L = 1 - cos(h_E^{l_t}(P_gen), sg[h_E^{l_s}(P_mem)])", "size": 11, "italic": True},
    {"text": "Minimizes cosine distance between reasoning-layer representation and storage representation. Stop-gradient on storage (teacher).", "size": 11},
    {"text": ""},
    {"text": "ContraRoute (InfoNCE Contrastive Routing):", "size": 12, "bold": True},
    {"text": "  L = -log[ exp(sim(q,k+)/τ) / (exp(sim(q,k+)/τ) + Σ exp(sim(q,k-)/τ)) ]", "size": 11, "italic": True},
    {"text": "Pulls reasoning-layer rep toward same-fact storage rep; pushes away other facts in batch. τ=0.07.", "size": 11},
    {"text": ""},
    {"text": "Probe Loss:", "size": 12, "bold": True},
    {"text": "  L = CE(φ*(h_E^{l_t}(P_gen)), y*)", "size": 11, "italic": True},
    {"text": "Frozen linear probe φ* (trained on storage-layer reps) grades the reasoning layer. Forces mid-layer into linearly decodable subspace (PoLM 3.1 condition).", "size": 11},
    {"text": ""},
    {"text": "Hybrid: α·L_Probe + (1-α)·L_Contra, α=0.5", "size": 12, "bold": True},
    {"text": ""},
    {"text": "All 4 variants tested across 4 models × 2 datasets.", "size": 12, "italic": True},
])
print("✓ Slide 18 [LAST]: AA-SFT Methodology")

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
OUTPUT = "AT&T_CDO_DataScienceShow_SFT_KUG_v2.pptx"
prs.save(OUTPUT)
print(f"\n✅  Saved: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides):
    title_shape = sl.shapes.title
    t = (title_shape.text if title_shape else "(no title)").replace("\n", " | ")
    print(f"   [{i+1:02d}] {t[:75]}")
