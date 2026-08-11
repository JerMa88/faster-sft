"""
Build: AT&T CDO Data Science Show Presentation
"Bridging the Knowing-Using Gap in Enterprise AI"
~30-minute talk for Product Owners, Investors, Business Personnel.

Template: slide_format.pptx
Output:  AT&T_CDO_DataScienceShow_SFT_KUG.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import copy
import lxml.etree as etree

# ──────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────

FOOTER_TEXT = "AT&T CDO Data Science Show / August 2026 / © 2026 AT&T Intellectual Property - AT&T Proprietary (Internal)"
ATT_BLUE = RGBColor(0x00, 0x9F, 0xDB)   # AT&T signature blue


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"Layout not found: {name}")


def add_slide(prs, layout_name):
    layout = get_layout(prs, layout_name)
    slide = prs.slides.add_slide(layout)
    return slide


def set_ph(slide, idx, text, bold=False, font_size=None, color=None):
    """Set text in placeholder by index. Clears and rewrites."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            if bold:
                run.font.bold = True
            if font_size:
                run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            return ph
    # placeholder not found – silently skip
    return None


def set_ph_multiline(slide, idx, lines, bold_first=False, font_size=None, color=None, bullet=False):
    """Set multi-line text in placeholder (list of strings → paragraphs)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = line
                if i == 0 and bold_first:
                    run.font.bold = True
                if font_size:
                    run.font.size = Pt(font_size)
                if color:
                    run.font.color.rgb = color
            return ph
    return None


def set_footer(slide, text=FOOTER_TEXT):
    set_ph(slide, 12, text)


def set_slide_number(slide, num):
    set_ph(slide, 11, str(num))


# ──────────────────────────────────────────────────────────────
# Load template & strip example slides (keep masters/layouts)
# ──────────────────────────────────────────────────────────────

prs = Presentation("slide_format.pptx")

# Remove all existing slides while preserving slide layouts
# Use the correct XML namespace prefix for the r:id attribute
from pptx.oxml.ns import qn

while len(prs.slides) > 0:
    slide_elem = prs.slides._sldIdLst[0]
    # The relationship id attribute uses the 'r' namespace
    rId = slide_elem.get(qn('r:id'))
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
    del prs.slides._sldIdLst[0]

print(f"Cleared slides. Layouts available: {len(prs.slide_layouts)}")


# ──────────────────────────────────────────────────────────────
# SLIDE 1: Title / Cover
# Layout: Cover w/ Subtitle (White)
# ──────────────────────────────────────────────────────────────
slide1 = add_slide(prs, "Cover w/ Subtitle (White)")
set_ph(slide1, 0,
       "Bridging the Knowing-Using Gap in Enterprise AI",
       bold=True)
set_ph(slide1, 13,
       "Evolution of SFT, Architectural Insights & Next-Gen Model Adaptation")
# Additional lines via idx 10 / 11 for author/org
set_ph(slide1, 10, "AT&T Chief Data Office — AI Research Initiative")
set_ph(slide1, 11, "CDO Data Science Show  |  August 2026")
print("Slide 1 done: Cover")

# ──────────────────────────────────────────────────────────────
# SLIDE 2: Executive Summary
# Layout: Title + Subtitle (2 column w/ headers)
# ──────────────────────────────────────────────────────────────
slide2 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(slide2, 0, "Today's Talk: Two Questions, One Solution")
set_ph(slide2, 13, "A 30-minute journey from fine-tuning fundamentals to AT&T's AI research frontier")
set_footer(slide2)

# Left column header
set_ph(slide2, 14, "The Business Promise of LLM Fine-Tuning")
# Left column body
set_ph_multiline(slide2, 18, [
    "Adapting powerful open-source AI models to AT&T's domain — network, billing, customer service — via Supervised Fine-Tuning (SFT) is fast, private, and cost-effective.",
    "",
    "The expectation: a model trained on AT&T facts should be able to reason about those facts to solve complex multi-step customer and operational queries.",
])

# Right column header
set_ph(slide2, 15, "The Hidden Bottleneck — and Our Solution")
# Right column body
set_ph_multiline(slide2, 19, [
    "Standard SFT creates models that memorize facts perfectly but fail to use them in multi-step logic — a failure mode we call the Knowing-Using Gap (KUG).",
    "",
    "CDO Research: We've identified the architectural root cause and built a fix that closes this gap at training time — with zero added deployment cost or latency.",
])
print("Slide 2 done: Executive Summary")

# ──────────────────────────────────────────────────────────────
# SLIDE 3: Section Divider — Part I
# Layout: Divider 01
# ──────────────────────────────────────────────────────────────
slide3 = add_slide(prs, "Divider 01")
set_ph(slide3, 0, "PART I\nSupervised Fine-Tuning (SFT)\nin the Enterprise")
set_footer(slide3)
print("Slide 3 done: Divider I")

# ──────────────────────────────────────────────────────────────
# SLIDE 4: What is SFT?
# Layout: Title + Subtitle (3 column w/ headers)
# ──────────────────────────────────────────────────────────────
slide4 = add_slide(prs, "Title + Subtitle (3 column w/ headers)")
set_ph(slide4, 0, "What Is SFT — and Why Does AT&T Depend On It?")
set_ph(slide4, 13, "Supervised Fine-Tuning is the industry-standard technique for teaching AI models domain-specific knowledge")
set_footer(slide4)

# Col 1
set_ph(slide4, 14, "💰  Cost-Efficient Adaptation")
set_ph_multiline(slide4, 18, [
    "Adapts billion-parameter models using Low-Rank Adaptation (LoRA) — a thin layer of trainable adapters — avoiding multi-million dollar full training runs.",
    "",
    "Runs on AT&T's existing GPU infrastructure. Fast turnaround: days, not months.",
])

# Col 2
set_ph(slide4, 15, "🔒  Data Sovereignty & Privacy")
set_ph_multiline(slide4, 19, [
    "Fine-tuning happens in-house. Proprietary AT&T network specs, customer data, and billing logic stay inside corporate firewalls — no third-party API exposure.",
    "",
    "Satisfies regulatory and compliance requirements for sensitive enterprise data.",
])

# Col 3
set_ph(slide4, 20, "🎯  Domain Customization at Scale")
set_ph_multiline(slide4, 21, [
    "Injects AT&T-specific vocabulary, diagnostic workflows, billing rules, and enterprise knowledge graph facts directly into model parameters.",
    "",
    "Result: a domain-expert AI that outperforms generic GPT-4-style APIs on AT&T-specific tasks.",
])
print("Slide 4 done: What is SFT?")

# ──────────────────────────────────────────────────────────────
# SLIDE 5: Section Divider — Part II
# Layout: Divider 01
# ──────────────────────────────────────────────────────────────
slide5 = add_slide(prs, "Divider 01")
set_ph(slide5, 0, 'PART II\nThe \u201cKnowing-Using Gap\u201d (KUG)\nWhy SFT Fails Business Multi-Step Logic')
set_footer(slide5)
print("Slide 5 done: Divider II")

# ──────────────────────────────────────────────────────────────
# SLIDE 6: The KUG Explained
# Layout: Large Text Block w/ Blue Curve
# ──────────────────────────────────────────────────────────────
slide6 = add_slide(prs, "Large Text Block w/ Blue Curve")
set_ph(slide6, 0, "The Knowing-Using Gap (KUG)")
set_ph(slide6, 14,
       '"Fine-tuned AI aces the recall test but fails the reasoning execution test — every single time."')
set_footer(slide6)

# Three column items: col header / col body
set_ph(slide6, 15, "✅  What Succeeds: Memorization")
set_ph_multiline(slide6, 18, [
    'Query: "What is the max speed of Fiber Tier 500?"',
    "AI Output: \"500 Mbps\" → Correct.",
    "",
    "The model perfectly recalls the fact it was trained on — like a student who memorized a flashcard.",
])

set_ph(slide6, 19, "❌  What Fails: Multi-Hop Reasoning")
set_ph_multiline(slide6, 20, [
    'Query: "Does Router Model X support the max speed of Fiber Tier 500?"',
    "AI Output: Fails or Hallucinates → ~0% Success.",
    "",
    "The model cannot connect two known facts into a reasoned answer — the operational requirement that actually matters.",
])

set_ph(slide6, 21, "📊  Empirical Scale of the Problem")
set_ph_multiline(slide6, 22, [
    "In AT&T-style knowledge benchmark experiments across 4 AI models:",
    "• Memorization accuracy (Amem): 24% – 74%",
    "• Reasoning accuracy (Agen):  0.9% – 25%",
    "• KUG Ratio: up to 65× gap between knowing and using.",
])
print("Slide 6 done: KUG Explained")

# ──────────────────────────────────────────────────────────────
# SLIDE 7: Inside the LLM — Storage vs Reasoning Circuits
# Layout: Title + Subtitle (2 column w/ headers)
# ──────────────────────────────────────────────────────────────
slide7 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(slide7, 0, "Why This Happens: Inside the Brain of an LLM")
set_ph(slide7, 13, "Transformer AI models have distinct internal zones — and standard training short-circuits them")
set_footer(slide7)

set_ph(slide7, 14, "🗄️  The Filing Cabinet (Early & Late Layers)")
set_ph_multiline(slide7, 18, [
    "Standard fine-tuning writes new facts into the model's early and late layers — acting as filing cabinets for factual associations.",
    "",
    "The optimization loss reaches ZERO here. Fine-tuning stops. The system declares: 'Learning complete!'",
    "",
    "The fact is stored — but stored in the wrong place for business reasoning.",
])

set_ph(slide7, 15, "🧠  The Executive Desk (Middle Layers) — Disconnected")
set_ph_multiline(slide7, 19, [
    "Multi-hop reasoning — connecting facts, chaining logic, cross-referencing — happens in the middle layers. Think of these as the executive desks where strategy is made.",
    "",
    "The problem: standard SFT never builds a conveyor belt from the filing cabinet to the executive desk.",
    "",
    "Academically confirmed: new facts are stored in layers 1-8, but reasoning circuits are in layers 10-20. No highway was built between them.",
])
print("Slide 7 done: Storage vs Reasoning Circuits")

# ──────────────────────────────────────────────────────────────
# SLIDE 8: Section Divider — Part III
# Layout: Divider 01
# ──────────────────────────────────────────────────────────────
slide8 = add_slide(prs, "Divider 01")
set_ph(slide8, 0, "PART III\nThe Landscape of SFT Solutions\nWhere the Industry Stands Today")
set_footer(slide8)
print("Slide 8 done: Divider III")

# ──────────────────────────────────────────────────────────────
# SLIDE 9: Current SFT Approaches & Their Limits
# Layout: Title + Subtitle (4 column w/ headers)
# ──────────────────────────────────────────────────────────────
slide9 = add_slide(prs, "Title + Subtitle (4 column w/ headers)")
set_ph(slide9, 0, "State of the Art: Four Approaches & Their Limitations")
set_ph(slide9, 13, "Every existing solution has a critical gap that makes it impractical for AT&T enterprise deployments")
set_footer(slide9)

# Col 1: Standard SFT
set_ph(slide9, 14, "1. Standard SFT / LoRA")
set_ph_multiline(slide9, 18, [
    "The industry default.",
    "",
    "✅ Fast, cheap, efficient.",
    "✅ Easy to deploy.",
    "",
    "❌ Inherits the full KUG — facts never reach reasoning circuits.",
    "❌ Multi-hop accuracy near zero.",
])

# Col 2: Model Editing
set_ph(slide9, 15, "2. Model Editing (ROME/MEMIT)")
set_ph_multiline(slide9, 19, [
    "Surgically modifies specific model weights to insert facts.",
    "",
    "✅ Targeted & fast for single facts.",
    "",
    "❌ 'Hopping-too-late' failure: edited facts are placed in layers the reasoning engine can't access for multi-step queries.",
    "❌ Breaks for chains of 2+ facts.",
])

# Col 3: Knowledge Distillation
set_ph(slide9, 20, "3. Knowledge Distillation")
set_ph_multiline(slide9, 21, [
    "Trains a small model by mimicking a larger teacher model's outputs.",
    "",
    "✅ Creates compact, fast models.",
    "",
    "❌ Requires two full models running simultaneously — expensive.",
    "❌ Doesn't address the internal routing failure within a single model.",
])

# Col 4: Diagnostic Patching
set_ph(slide9, 22, "4. Diagnostic Self-Patching")
set_ph_multiline(slide9, 23, [
    "Manually copies hidden state vectors to middle layers at test time. Proves 58-75% of lost reasoning is recoverable.",
    "",
    "✅ Proves the problem is fixable.",
    "",
    "❌ Requires knowing the answer ahead of time.",
    "❌ Slows down every inference response.",
    "❌ Not deployable in production.",
])
print("Slide 9 done: Existing Approaches")

# ──────────────────────────────────────────────────────────────
# SLIDE 10: Section Divider — Part IV
# Layout: Divider 01
# ──────────────────────────────────────────────────────────────
slide10 = add_slide(prs, "Divider 01")
set_ph(slide10, 0, "PART IV\nEvolving Enterprise SFT\nStrategic Opportunity & Business Value")
set_footer(slide10)
print("Slide 10 done: Divider IV")

# ──────────────────────────────────────────────────────────────
# SLIDE 11: Evolution of Fine-Tuning
# Layout: Title + Subtitle (1/2 Blue Curve)
# ──────────────────────────────────────────────────────────────
slide11 = add_slide(prs, "Title + Subtitle (1/2 Blue Curve)")
set_ph(slide11, 0, "Evolving Fine-Tuning: Phase 3 Has Arrived")
# Big text block on left (idx 14)
set_ph_multiline(slide11, 14, [
    "The AI industry is entering a new era of model adaptation — moving from text-guessing to concept-routing.",
    "",
    "AT&T CDO is at the frontier of Phase 3.",
])
set_footer(slide11)

# 4 bullet items on right (idx 22/23/24/25/26/27/28/29)
set_ph(slide11, 22, "Phase 1 (Pre-2022): Full Fine-Tuning (FFT)")
set_ph_multiline(slide11, 23, [
    "Train all model parameters. High compute cost. Risk of forgetting old knowledge. Inflexible for rapid updates.",
])

set_ph(slide11, 24, "Phase 2 (2022–2024): Parameter-Efficient SFT (LoRA)")
set_ph_multiline(slide11, 25, [
    "Train only small adapter layers. Efficient and practical. Industry standard today — but the KUG persists.",
])

set_ph(slide11, 26, "Phase 3 (2025+): Alignment-Aware Fine-Tuning")
set_ph_multiline(slide11, 27, [
    "Guide the model's internal representation pathways during training. Fix the routing failure at its source. No added runtime cost.",
])

set_ph(slide11, 28, "Where AT&T CDO Research Is Operating →")
set_ph_multiline(slide11, 29, [
    "Building Phase 3 techniques tailored to enterprise knowledge graphs: customer care, network operations, billing intelligence.",
])
print("Slide 11 done: Evolution of Fine-Tuning")

# ──────────────────────────────────────────────────────────────
# SLIDE 12: Strategic ROI & Business Value
# Layout: Title + Subtitle (3 column w/ headers)
# ──────────────────────────────────────────────────────────────
slide12 = add_slide(prs, "Title + Subtitle (3 column w/ headers)")
set_ph(slide12, 0, "Strategic Value: Why This Matters for AT&T")
set_ph(slide12, 13, "Closing the Knowing-Using Gap delivers compounding business value with zero added deployment cost")
set_footer(slide12)

# Col 1
set_ph(slide12, 14, "📈  Higher-Accuracy AI Decisions")
set_ph_multiline(slide12, 18, [
    "Moves AI from single-fact lookup to genuine multi-step reasoning.",
    "",
    "Customer service bots correctly cross-reference plan features and device compatibility.",
    "",
    "Network operations assistants chain alerts to root-cause diagnoses without human re-prompting.",
])

# Col 2
set_ph(slide12, 15, "⚡  Zero Inference Latency Overhead")
set_ph_multiline(slide12, 19, [
    "The routing fix is baked in at training time.",
    "",
    "Once fine-tuning is complete, the model runs at exactly the same speed as standard SFT.",
    "",
    "No extra inference hooks, no second model in memory, no answer lookup required.",
    "",
    "Identical deployment footprint to today's LoRA-SFT pipelines.",
])

# Col 3
set_ph(slide12, 20, "🔌  Plug-and-Play Integration")
set_ph_multiline(slide12, 21, [
    "Integrates directly into existing AT&T LoRA fine-tuning infrastructure.",
    "",
    "No hardware upgrades required — runs on existing A100 GPU clusters.",
    "",
    "Modular: applies to any enterprise knowledge domain — customer plans, network topology, internal policy.",
])
print("Slide 12 done: Business Value")

# ──────────────────────────────────────────────────────────────
# SLIDE 13: Section Divider — Part V
# Layout: Divider 01
# ──────────────────────────────────────────────────────────────
slide13 = add_slide(prs, "Divider 01")
set_ph(slide13, 0, "PART V\nAT&T CDO Researcher Spotlight\nAlignment-Aware SFT — Methodology")
set_footer(slide13)
print("Slide 13 done: Divider V")

# ──────────────────────────────────────────────────────────────
# SLIDE 14 [LAST]: Researcher Spotlight — AA-SFT Methodology
# Layout: Title + Subtitle (2 column w/ headers)
# ──────────────────────────────────────────────────────────────
slide14 = add_slide(prs, "Title + Subtitle (2 column w/ headers)")
set_ph(slide14, 0, "Initiative Spotlight: Alignment-Aware SFT (AA-SFT)")
set_ph(slide14, 13,
       "A novel training-time auxiliary loss that builds the internal routing highway — permanently, during fine-tuning, at zero inference cost")
set_footer(slide14)

# Left column: The Dual-Forward Pass Framework
set_ph(slide14, 14, "🔁  The Dual-Forward Pass Training Framework")
set_ph_multiline(slide14, 18, [
    "For every enterprise fact in the training data, two prompts are used simultaneously:",
    "",
    "1. Memorization Prompt (P_mem): A direct recall question.",
    "   → Captures where the model stores the fact (early/late storage layers).",
    "",
    "2. Reasoning Prompt (P_gen): A multi-hop inference question.",
    "   → Captures what the model does with the fact in reasoning layers.",
    "",
    "An Alignment Loss (L_align) is added to standard SFT loss:",
    "   L_total = L_SFT + λ · L_align",
    "",
    "This forces the model's middle (reasoning) layers under P_gen to match the storage representations from P_mem — building the internal routing highway during training.",
    "",
    "At deployment: standard single forward pass. Zero overhead.",
])

# Right column: Four Loss Variants
set_ph(slide14, 15, "⚙️  Four Alignment Loss Variants (Research Exploration)")
set_ph_multiline(slide14, 19, [
    "RepDist (Representation Distillation):",
    "Minimizes cosine distance between reasoning-layer representation and stored fact vector. Like teaching the middle layers to 'point in the same direction' as the filing cabinet.",
    "",
    "ContraRoute (Contrastive Routing):",
    "Uses contrastive learning — pulls reasoning representations toward the correct fact while pushing them away from other distractors in the batch. Sharpens fact discrimination.",
    "",
    "Probe Loss:",
    "Trains a frozen 'fact decoder' on stored representations, then uses it to grade the reasoning layer. Forces reasoning layers into a decodable knowledge subspace.",
    "",
    "Hybrid:",
    "Combines Probe and ContraRoute losses with equal weighting for complementary alignment signals.",
    "",
    "Tested across: 4 AI models (1.2B–4B parameters) × 2 enterprise knowledge benchmarks.",
])
print("Slide 14 done: AA-SFT Methodology [LAST SLIDE]")

# ──────────────────────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────────────────────
OUTPUT = "AT&T_CDO_DataScienceShow_SFT_KUG.pptx"
prs.save(OUTPUT)
print(f"\n✅  Saved: {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides):
    title_shape = sl.shapes.title
    t = title_shape.text if title_shape else "(no title)"
    print(f"   [{i+1:02d}] {t[:70]}")
